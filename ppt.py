import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from langchain.llms import OpenAI
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate


# 创建翻译链
def create_translation_chain():
    llm = OpenAI(temperature=0.7)
    template = """
    将以下文本翻译成英文：

    {text}

    翻译：
    """
    prompt = PromptTemplate(input_variables=["text"], template=template)
    return LLMChain(llm=llm, prompt=prompt)

# 翻译文本框内容并保持格式
def translate_text_frame(text_frame, translation_chain):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                original_text = run.text
                translated_text = translation_chain.run(text=original_text).strip()
                
                # 存储原始格式
                font_name = run.font.name
                font_size = run.font.size
                font_bold = run.font.bold
                font_italic = run.font.italic
                font_color = None
                
                # 设置翻译后的文本
                run.text = translated_text
                
                # 应用原始格式
                run.font.name = font_name
                run.font.size = font_size
                run.font.bold = font_bold
                run.font.italic = font_italic
                if font_color:
                    run.font.color.rgb = font_color
                
                print(f"Translated: '{original_text}' to '{translated_text}'")

# 处理幻灯片中的所有形状
def process_shapes(shapes, translation_chain):
    for shape in shapes:
        if shape.has_text_frame:
            translate_text_frame(shape.text_frame, translation_chain)
        elif hasattr(shape, "shapes"):  # 处理组内形状
            process_shapes(shape.shapes, translation_chain)

# 替换文本并保持格式
def replace_text_in_pptx(input_pptx_path, output_pptx_path, progress_callback):
    prs = Presentation(input_pptx_path)
    translation_chain = create_translation_chain()

    total_shapes = sum([len(slide.shapes) for slide in prs.slides])
    processed_shapes = 0

    for slide in prs.slides:
        process_shapes(slide.shapes, translation_chain)
        processed_shapes += len(slide.shapes)
        progress = processed_shapes / total_shapes * 100
        progress_callback(progress)

    prs.save(output_pptx_path)
    progress_callback(100)  # 完成时确保进度设置为100%

def main():
    input_pptx_path = 'file/autobot.pptx'
    output_pptx_path = 'file/output_translated.pptx'
    
    def dummy_progress_callback(progress):
        print(f"Progress: {progress}%")
    
    replace_text_in_pptx(input_pptx_path, output_pptx_path, dummy_progress_callback)
    print(f"PPT document saved as {output_pptx_path} with translated text")

if __name__ == '__main__':
    main()
